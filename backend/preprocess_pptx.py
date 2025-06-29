#!/usr/bin/env python3
"""
PPTX Preprocessing Script for Gym Chatbot Knowledge Base

This script processes PowerPoint (.pptx) files to extract content, tables, and color coding,
then formats the output for ingestion into the AWS Bedrock Knowledge Base.

Features:
- Extracts text content from slides
- Processes tables with color coding interpretation
- Handles slide layouts and formatting
- Converts to structured text format for knowledge base
- Supports batch processing of multiple files
"""

import os
import sys
import json
import argparse
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from dataclasses import dataclass
import logging

# Add the current directory to Python path for imports
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("❌ Required library 'python-pptx' not found.")
    print("Install it with: pip install python-pptx")
    sys.exit(1)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('pptx_preprocessing.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class TableData:
    """Represents processed table data with color coding information"""
    headers: List[str]
    rows: List[List[str]]
    color_coding: Dict[str, str]  # cell_position -> color_meaning
    slide_number: int
    table_index: int

@dataclass
class SlideContent:
    """Represents processed slide content"""
    slide_number: int
    title: str
    text_content: List[str]
    tables: List[TableData]
    notes: str

class ColorCodingInterpreter:
    """Interprets color coding in tables and shapes"""
    
    # Common color meanings in fitness/gym context
    COLOR_MEANINGS = {
        # Green variations - typically positive/good
        '00FF00': 'Excellent/Great',
        '008000': 'Good',
        '90EE90': 'Good',
        '32CD32': 'Good',
        '228B22': 'Good',
        
        # Yellow variations - typically moderate/average
        'FFFF00': 'Moderate/Average',
        'FFD700': 'Moderate/Average',
        'FFA500': 'Moderate/Average',
        'FF8C00': 'Moderate/Average',
        
        # Red variations - typically poor/needs attention
        'FF0000': 'Poor/Needs Attention',
        'DC143C': 'Poor/Needs Attention',
        'B22222': 'Poor/Needs Attention',
        '8B0000': 'Poor/Needs Attention',
        
        # Blue variations - typically informational
        '0000FF': 'Informational',
        '4169E1': 'Informational',
        '1E90FF': 'Informational',
        '00BFFF': 'Informational',
        
        # Gray variations - typically neutral
        '808080': 'Neutral',
        'A9A9A9': 'Neutral',
        'C0C0C0': 'Neutral',
        'D3D3D3': 'Neutral',
    }
    
    @staticmethod
    def rgb_to_hex(rgb_color: RGBColor) -> Optional[str]:
        """Convert RGB color to hex string"""
        if rgb_color is None:
            return None
        return f"{rgb_color.r:02X}{rgb_color.g:02X}{rgb_color.b:02X}"
    
    @staticmethod
    def interpret_color(color_hex: str) -> str:
        """Interpret color meaning based on hex value"""
        if not color_hex:
            return "Default"
        
        # Validate hex color format
        if not isinstance(color_hex, str) or len(color_hex) != 6:
            logger.warning(f"Invalid hex color format: {color_hex}")
            return "Default"
        
        try:
            # Validate that it's a valid hex string
            int(color_hex, 16)
        except ValueError:
            logger.warning(f"Invalid hex color value: {color_hex}")
            return "Default"
        
        # Exact match
        if color_hex.upper() in ColorCodingInterpreter.COLOR_MEANINGS:
            return ColorCodingInterpreter.COLOR_MEANINGS[color_hex.upper()]
        
        # Find closest match (simple distance calculation)
        best_match = "Default"
        min_distance = float('inf')
        
        for hex_color, meaning in ColorCodingInterpreter.COLOR_MEANINGS.items():
            try:
                # Convert hex to RGB for distance calculation
                target_r = int(hex_color[0:2], 16)
                target_g = int(hex_color[2:4], 16)
                target_b = int(hex_color[4:6], 16)
                
                current_r = int(color_hex[0:2], 16)
                current_g = int(color_hex[2:4], 16)
                current_b = int(color_hex[4:6], 16)
                
                # Calculate Euclidean distance
                distance = ((target_r - current_r) ** 2 + 
                           (target_g - current_g) ** 2 + 
                           (target_b - current_b) ** 2) ** 0.5
                
                if distance < min_distance:
                    min_distance = distance
                    best_match = meaning
            except (ValueError, IndexError) as e:
                logger.warning(f"Error calculating color distance for {hex_color}: {e}")
                continue
        
        return best_match

class PPTXProcessor:
    """Main class for processing PPTX files"""
    
    def __init__(self, output_format: str = 'structured'):
        self.output_format = output_format
        self.color_interpreter = ColorCodingInterpreter()
    
    def process_pptx_file(self, file_path: str) -> Dict[str, Any]:
        """Process a single PPTX file and return structured data"""
        try:
            logger.info(f"Processing PPTX file: {file_path}")
            
            # Load presentation
            prs = Presentation(file_path)
            
            # Extract metadata
            metadata = self._extract_metadata(prs, file_path)
            
            # Process slides
            slides_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = self._process_slide(slide, slide_num)
                slides_content.append(slide_content)
            
            # Generate output
            output = {
                'metadata': metadata,
                'slides': slides_content,
                'summary': self._generate_summary(slides_content)
            }
            
            logger.info(f"Successfully processed {len(slides_content)} slides")
            return output
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            raise
    
    def _extract_metadata(self, prs: Presentation, file_path: str) -> Dict[str, Any]:
        """Extract presentation metadata"""
        metadata = {
            'filename': Path(file_path).name,
            'file_path': file_path,
            'total_slides': len(prs.slides),
            'slide_size': {
                'width': prs.slide_width,
                'height': prs.slide_height
            }
        }
        
        # Extract core properties if available
        if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
            metadata['title'] = prs.core_properties.title
        if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
            metadata['author'] = prs.core_properties.author
        if hasattr(prs.core_properties, 'subject') and prs.core_properties.subject:
            metadata['subject'] = prs.core_properties.subject
        
        return metadata
    
    def _process_slide(self, slide, slide_number: int) -> SlideContent:
        """Process a single slide and extract all content"""
        title = ""
        text_content = []
        tables = []
        notes = ""
        
        # Extract slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        
        # Process shapes in the slide
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = self._process_table(shape, slide_number, len(tables))
                tables.append(table_data)
            elif hasattr(shape, 'text_frame'):
                text = shape.text_frame.text.strip()
                if text:
                    if not title and shape == slide.shapes[0]:  # First shape is usually title
                        title = text
                    else:
                        text_content.append(text)
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            tables=tables,
            notes=notes
        )
    
    def _process_table(self, shape, slide_number: int, table_index: int) -> TableData:
        """Process a table shape and extract data with color coding"""
        try:
            table = shape.table
            headers = []
            rows = []
            color_coding = {}
            
            # Safely get table dimensions
            try:
                num_rows = len(table.rows)
                num_cols = len(table.rows[0].cells) if table.rows else 0
            except (IndexError, AttributeError):
                logger.debug(f"Could not determine table dimensions on slide {slide_number}")
                return TableData(
                    headers=[],
                    rows=[],
                    color_coding={},
                    slide_number=slide_number,
                    table_index=table_index
                )
            
            # Process headers (first row)
            if num_rows > 0:
                try:
                    for cell in table.rows[0].cells:
                        headers.append(cell.text.strip())
                except Exception as e:
                    logger.debug(f"Error processing headers on slide {slide_number}: {e}")
                    headers = []
            
            # Process data rows
            for row_idx in range(1, num_rows):
                try:
                    row = table.rows[row_idx]
                    row_data = []
                    
                    for col_idx in range(num_cols):
                        try:
                            if col_idx < len(row.cells):
                                cell = row.cells[col_idx]
                                cell_text = cell.text.strip()
                                row_data.append(cell_text)
                                
                                # Analyze cell color (skip if there are issues)
                                try:
                                    cell_color = self._get_cell_color(cell)
                                    if cell_color:
                                        position = f"{row_idx},{col_idx}"
                                        color_meaning = self.color_interpreter.interpret_color(cell_color)
                                        color_coding[position] = color_meaning
                                except Exception as e:
                                    logger.debug(f"Error processing cell color at {row_idx},{col_idx}: {e}")
                            else:
                                row_data.append("")
                        except Exception as e:
                            logger.debug(f"Error processing cell at {row_idx},{col_idx}: {e}")
                            row_data.append("")
                    
                    rows.append(row_data)
                except Exception as e:
                    logger.debug(f"Error processing row {row_idx} on slide {slide_number}: {e}")
                    continue
            
            return TableData(
                headers=headers,
                rows=rows,
                color_coding=color_coding,
                slide_number=slide_number,
                table_index=table_index
            )
        except Exception as e:
            logger.warning(f"Error processing table on slide {slide_number}: {e}")
            # Return empty table data instead of crashing
            return TableData(
                headers=[],
                rows=[],
                color_coding={},
                slide_number=slide_number,
                table_index=table_index
            )
    
    def _get_cell_color(self, cell) -> Optional[str]:
        """Extract background color from a table cell"""
        try:
            # Check for fill color
            if hasattr(cell, 'fill') and cell.fill.type:
                if hasattr(cell.fill.fore_color, 'rgb') and cell.fill.fore_color.rgb:
                    color_hex = self.color_interpreter.rgb_to_hex(cell.fill.fore_color.rgb)
                    if color_hex and len(color_hex) == 6:
                        return color_hex
                    else:
                        logger.debug(f"Invalid color hex extracted: {color_hex}")
        except Exception as e:
            logger.debug(f"Error extracting cell color: {e}")
        
        return None
    
    def _generate_summary(self, slides_content: List[SlideContent]) -> Dict[str, Any]:
        """Generate a summary of the presentation content"""
        total_tables = sum(len(slide.tables) for slide in slides_content)
        total_text_blocks = sum(len(slide.text_content) for slide in slides_content)
        
        # Count color-coded cells
        total_color_coded_cells = 0
        color_distribution = {}
        
        for slide in slides_content:
            for table in slide.tables:
                total_color_coded_cells += len(table.color_coding)
                for color_meaning in table.color_coding.values():
                    color_distribution[color_meaning] = color_distribution.get(color_meaning, 0) + 1
        
        return {
            'total_slides': len(slides_content),
            'total_tables': total_tables,
            'total_text_blocks': total_text_blocks,
            'total_color_coded_cells': total_color_coded_cells,
            'color_distribution': color_distribution
        }
    
    def format_for_knowledge_base(self, processed_data: Dict[str, Any]) -> str:
        """Format processed data for knowledge base ingestion"""
        output_lines = []
        
        # Add metadata
        metadata = processed_data['metadata']
        output_lines.append(f"# {metadata.get('title', metadata['filename'])}")
        output_lines.append("")
        output_lines.append(f"**Source:** {metadata['filename']}")
        if 'author' in metadata:
            output_lines.append(f"**Author:** {metadata['author']}")
        if 'subject' in metadata:
            output_lines.append(f"**Subject:** {metadata['subject']}")
        output_lines.append(f"**Total Slides:** {metadata['total_slides']}")
        output_lines.append("")
        
        # Process each slide
        for slide in processed_data['slides']:
            # Add slide header
            if slide.title.strip():
                output_lines.append(f"## Slide {slide.slide_number}: {slide.title}")
            else:
                output_lines.append(f"## Slide {slide.slide_number}")
            output_lines.append("")
            
            # Add text content with proper formatting
            for text in slide.text_content:
                if text.strip():
                    # Split text into paragraphs and format properly
                    paragraphs = text.split('\n')
                    for paragraph in paragraphs:
                        paragraph = paragraph.strip()
                        if paragraph:
                            # Check if it looks like a list item
                            if paragraph.startswith(('•', '-', '*', '1)', '2)', '3)', '4)', '5)', '6)', '7)', '8)', '9)', '0)')):
                                output_lines.append(paragraph)
                            elif paragraph.startswith(('A)', 'B)', 'C)', 'D)', 'E)', 'F)', 'G)', 'H)')):
                                output_lines.append(paragraph)
                            elif paragraph.startswith(('WEEK', 'DAY', 'CONDITIONING', 'GUÍA', 'INTENSITY', 'REST')):
                                output_lines.append(f"**{paragraph}**")
                            elif paragraph.startswith(('*', '_')) and paragraph.endswith(('*', '_')):
                                # Italic text
                                output_lines.append(paragraph)
                            else:
                                # Regular paragraph
                                output_lines.append(paragraph)
                    output_lines.append("")
            
            # Add tables
            for table in slide.tables:
                if table.headers or table.rows:
                    output_lines.append(f"### Table {table.table_index + 1}")
                    output_lines.append("")
                    
                    # Add headers
                    if table.headers:
                        header_row = "| " + " | ".join(table.headers) + " |"
                        separator_row = "| " + " | ".join(["---"] * len(table.headers)) + " |"
                        output_lines.append(header_row)
                        output_lines.append(separator_row)
                    
                    # Add data rows with color coding
                    for row_idx, row in enumerate(table.rows, 1):
                        if row:  # Only add non-empty rows
                            formatted_row = []
                            for col_idx, cell in enumerate(row):
                                position = f"{row_idx},{col_idx}"
                                if position in table.color_coding:
                                    color_meaning = table.color_coding[position]
                                    formatted_cell = f"{cell} *({color_meaning})*"
                                else:
                                    formatted_cell = cell
                                formatted_row.append(formatted_cell)
                            row_text = "| " + " | ".join(formatted_row) + " |"
                            output_lines.append(row_text)
                    
                    output_lines.append("")
                    
                    # Add color coding legend if present
                    if table.color_coding:
                        output_lines.append("**Color Coding Legend:**")
                        unique_colors = set(table.color_coding.values())
                        for color in sorted(unique_colors):
                            output_lines.append(f"- {color}")
                        output_lines.append("")
            
            # Add slide notes if present
            if slide.notes and slide.notes.strip():
                output_lines.append("**Notes:**")
                # Format notes with proper paragraphs
                note_paragraphs = slide.notes.split('\n')
                for paragraph in note_paragraphs:
                    paragraph = paragraph.strip()
                    if paragraph:
                        output_lines.append(paragraph)
                output_lines.append("")
        
        # Add summary
        summary = processed_data['summary']
        output_lines.append("## Summary")
        output_lines.append("")
        output_lines.append(f"- **Total Slides:** {summary['total_slides']}")
        output_lines.append(f"- **Total Tables:** {summary['total_tables']}")
        output_lines.append(f"- **Total Text Blocks:** {summary['total_text_blocks']}")
        output_lines.append(f"- **Color-Coded Cells:** {summary['total_color_coded_cells']}")
        
        if summary['color_distribution']:
            output_lines.append("")
            output_lines.append("**Color Distribution:**")
            for color, count in summary['color_distribution'].items():
                output_lines.append(f"- {color}: {count} cells")
        
        return "\n".join(output_lines)

def process_pptx_files(input_paths: List[str], output_dir: str, format_type: str = 'structured') -> List[str]:
    """Process multiple PPTX files and save outputs"""
    processor = PPTXProcessor(output_format=format_type)
    output_files = []
    
    # Create output directory if it doesn't exist
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    for input_path in input_paths:
        try:
            # Process the file
            processed_data = processor.process_pptx_file(input_path)
            
            # Generate output filename
            input_filename = Path(input_path).stem
            output_filename = f"{input_filename}_processed.md"
            output_path = Path(output_dir) / output_filename
            
            # Format and save output
            if format_type == 'structured':
                # Save as structured JSON
                json_output_path = Path(output_dir) / f"{input_filename}_processed.json"
                with open(json_output_path, 'w', encoding='utf-8') as f:
                    json.dump(processed_data, f, indent=2, default=str)
                output_files.append(str(json_output_path))
            
            # Always save as markdown for knowledge base
            markdown_content = processor.format_for_knowledge_base(processed_data)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            output_files.append(str(output_path))
            logger.info(f"✅ Processed {input_path} -> {output_path}")
            
        except Exception as e:
            logger.error(f"❌ Failed to process {input_path}: {e}")
            continue
    
    return output_files

def main():
    """Main function for command-line usage"""
    parser = argparse.ArgumentParser(
        description="Preprocess PPTX files for gym chatbot knowledge base",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Process a single file
  python preprocess_pptx.py input.pptx -o output/

  # Process multiple files
  python preprocess_pptx.py *.pptx -o processed/

  # Process with structured JSON output
  python preprocess_pptx.py input.pptx -o output/ --format structured

  # Process directory of PPTX files
  python preprocess_pptx.py /path/to/pptx/files/ -o output/
        """
    )
    
    parser.add_argument(
        'input_paths',
        nargs='+',
        help='PPTX files or directories to process'
    )
    
    parser.add_argument(
        '-o', '--output',
        default='processed_pptx',
        help='Output directory for processed files (default: processed_pptx)'
    )
    
    parser.add_argument(
        '--format',
        choices=['structured', 'markdown'],
        default='structured',
        help='Output format: structured (JSON + MD) or markdown only (default: structured)'
    )
    
    parser.add_argument(
        '--verbose', '-v',
        action='store_true',
        help='Enable verbose logging'
    )
    
    args = parser.parse_args()
    
    if args.verbose:
        logger.setLevel(logging.DEBUG)
    
    # Expand input paths
    input_files = []
    for path in args.input_paths:
        path_obj = Path(path)
        if path_obj.is_file() and path_obj.suffix.lower() == '.pptx':
            input_files.append(str(path_obj))
        elif path_obj.is_dir():
            # Find all PPTX files in directory
            pptx_files = list(path_obj.glob('*.pptx'))
            input_files.extend([str(f) for f in pptx_files])
        else:
            logger.warning(f"Skipping {path}: not a PPTX file or directory")
    
    if not input_files:
        logger.error("No PPTX files found to process")
        sys.exit(1)
    
    logger.info(f"Found {len(input_files)} PPTX files to process")
    
    # Process files
    output_files = process_pptx_files(input_files, args.output, args.format)
    
    if output_files:
        logger.info(f"✅ Successfully processed {len(output_files)} files")
        logger.info(f"Output files saved to: {args.output}")
        for output_file in output_files:
            logger.info(f"  - {output_file}")
    else:
        logger.error("❌ No files were processed successfully")
        sys.exit(1)

if __name__ == "__main__":
    main() 